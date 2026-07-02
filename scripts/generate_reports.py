import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

REPORTS_DIR = "reports"
os.makedirs(REPORTS_DIR, exist_ok=True)

# ==========================================
# 1. Generate Final_Report.pdf
# ==========================================
def generate_pdf():
    pdf = FPDF()
    pdf.add_page()
    
    # Title
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, "Bluestock Fintech - Mutual Fund Capstone Project", ln=True, align='C')
    pdf.set_font("Arial", 'I', 12)
    pdf.cell(0, 10, "Final Analytics & Modeling Report", ln=True, align='C')
    pdf.ln(10)
    
    # Executive Summary
    pdf.set_font("Arial", 'B', 14)
    pdf.cell(0, 10, "1. Executive Summary", ln=True)
    pdf.set_font("Arial", '', 12)
    summary_text = (
        "This project comprehensively analyzes the performance of 40 Indian Mutual Fund schemes "
        "across the Large, Mid, and Small Cap categories over a 5-year period. An automated "
        "ETL pipeline ingested 17+ raw datasets into a SQLite star-schema database, enabling "
        "advanced quantitative performance modeling."
    )
    pdf.multi_cell(0, 7, summary_text)
    pdf.ln(5)
    
    # Methodology
    pdf.set_font("Arial", 'B', 14)
    pdf.cell(0, 10, "2. Methodology", ln=True)
    pdf.set_font("Arial", '', 12)
    methodology_text = (
        "- Calculated annualized CAGR (1yr, 3yr, 5yr) using 252 trading days.\n"
        "- Computed Sharpe & Sortino ratios against a 6.5% risk-free rate.\n"
        "- Analyzed CAPM metrics (Alpha, Beta) regressed against NIFTY 100.\n"
        "- Modeled downside risk through 95% Historical VaR and Maximum Drawdowns.\n"
    )
    pdf.multi_cell(0, 7, methodology_text)
    pdf.ln(5)
    
    # Key Findings
    pdf.set_font("Arial", 'B', 14)
    pdf.cell(0, 10, "3. Key Findings", ln=True)
    pdf.set_font("Arial", '', 12)
    findings_text = (
        "1. Risk & Return: Small Cap funds exhibited the highest 3-year absolute return, but "
        "endured the most severe maximum drawdowns (-22.3% on average).\n"
        "2. Portfolio Concentration: The Herfindahl-Hirschman Index (HHI) highlighted that IT "
        "and Pharma thematic funds possess severe diversification risks compared to standard Large-Cap indexes.\n"
        "3. Investor Behavior: The 2024 investor cohort accounts for over 40% of total inflows. "
        "Worryingly, 97.8% of frequent SIP investors demonstrate an average gap of over 35 days, "
        "indicating an 'at-risk' behavior pattern regarding SIP continuity."
    )
    pdf.multi_cell(0, 7, findings_text)
    
    pdf_path = os.path.join(REPORTS_DIR, "Final_Report.pdf")
    pdf.output(pdf_path)
    print(f"Generated PDF: {pdf_path}")

# ==========================================
# 2. Generate Presentation.pptx
# ==========================================
def generate_pptx():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Bluestock Fintech Capstone"
    subtitle.text = "Mutual Fund Data Analytics & Portfolio Modeling"
    
    # Executive Summary Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Executive Summary"
    tf = body_shape.text_frame
    tf.text = "Comprehensive analytics of 40 Mutual Fund schemes (5 Years)"
    p = tf.add_paragraph()
    p.text = "End-to-end Python ETL pipeline utilizing pandas & SQLite"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Advanced quantitative performance modeling (Sharpe, Alpha, VaR)"
    p.level = 1
    
    # Key Insights Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Top Analytical Insights"
    tf = body_shape.text_frame
    tf.text = "Small Cap Funds dominated returns but had ~22% Max Drawdowns."
    p = tf.add_paragraph()
    p.text = "High Sector HHI concentration found in Thematic funds."
    p = tf.add_paragraph()
    p.text = "97.8% of 'frequent' SIP investors have >35 day gap (high churn risk)."
    
    # Conclusion Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Final Recommendation"
    tf = body_shape.text_frame
    tf.text = "SBI Small Cap & Kotak Emerging Equity consistently outperformed."
    p = tf.add_paragraph()
    p.text = "Recommend continuous SIP to mitigate drawdown volatility."
    p.level = 1
    
    pptx_path = os.path.join(REPORTS_DIR, "Presentation.pptx")
    prs.save(pptx_path)
    print(f"Generated PPTX: {pptx_path}")

if __name__ == "__main__":
    generate_pdf()
    generate_pptx()
